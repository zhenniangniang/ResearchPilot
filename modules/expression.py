import io
import urllib.request
from utils.llm_client import chat
from utils.pdf_parser import extract_text_from_pdf, truncate_text
from prompts.expression_prompts import (
    EXPRESSION_SYSTEM,
    OUTLINE_USER,
    PPT_USER,
    PPT_FROM_FILE_USER,
    RESULT_PARAGRAPH_USER,
    CHART_DESCRIPTION_USER,
    MODEL_DIAGRAM_SYSTEM,
    MODEL_DIAGRAM_PROMPT_USER,
)


def generate_outline(topic: str, context: str = "") -> str:
    """生成开题报告提纲。"""
    user_prompt = OUTLINE_USER.format(topic=topic, context=context or "无")
    return chat(EXPRESSION_SYSTEM, user_prompt)


def generate_ppt_structure(topic: str, context: str = "", duration: int = 10) -> str:
    """生成 PPT 汇报结构框架（Markdown 文本）。"""
    user_prompt = PPT_USER.format(topic=topic, context=context or "无", duration=duration)
    return chat(EXPRESSION_SYSTEM, user_prompt)


def generate_ppt_from_file(file_bytes: bytes, filename: str, duration: int = 10) -> tuple[str, bytes]:
    """
    读取 PDF / Word 文件，AI 生成结构化 PPT 内容，并输出真实 .pptx 文件。
    返回 (markdown_preview, pptx_bytes)
    """
    # ── 1. 提取文字 ──────────────────────────────────────────────────────────
    fname = filename.lower()
    if fname.endswith(".pdf"):
        raw_text = extract_text_from_pdf(file_bytes)
    elif fname.endswith((".docx", ".doc")):
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        raw_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    else:
        raise ValueError(f"不支持的文件格式：{filename}，请上传 PDF 或 Word 文件")

    raw_text = truncate_text(raw_text)

    # ── 2. AI 生成结构化 PPT 内容 ─────────────────────────────────────────────
    user_prompt = PPT_FROM_FILE_USER.format(text=raw_text, duration=duration)
    ai_output = chat(EXPRESSION_SYSTEM, user_prompt)

    # ── 3. 解析 AI 输出并生成 .pptx ──────────────────────────────────────────
    pptx_bytes = _build_pptx(ai_output)
    return ai_output, pptx_bytes


def _build_pptx(ai_output: str) -> bytes:
    """将 AI 生成的结构化文本转换为美观的 .pptx 文件字节。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import re

    # ── 配色方案 ──────────────────────────────────────────────────────────────
    C_NAVY    = RGBColor(0x0F, 0x34, 0x60)   # 深蓝（标题栏背景）
    C_ACCENT  = RGBColor(0x00, 0xB4, 0xD8)   # 天蓝（装饰线、要点序号）
    C_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # 深夜蓝（封面背景）
    C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    C_LIGHT   = RGBColor(0xCC, 0xE5, 0xFF)   # 浅蓝（封面副文本）
    C_TEXT    = RGBColor(0x22, 0x22, 0x33)   # 正文深色
    C_GRAY    = RGBColor(0xAA, 0xAA, 0xBB)   # 页码灰
    C_BG      = RGBColor(0xF4, 0xF7, 0xFB)   # 内容页背景（浅灰蓝）

    prs = Presentation()
    W = Inches(13.33)
    H = Inches(7.5)
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    def _rgb_fill(shape, color: RGBColor):
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()   # 无边框

    def _add_textbox(slide, left, top, width, height,
                     text, font_size, bold=False, color=None,
                     align=PP_ALIGN.LEFT, word_wrap=True):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.size = Pt(font_size)
        p.font.bold = bold
        if color:
            p.font.color.rgb = color
        return tb

    def _add_title_slide(title: str, subtitle_lines: list[str]):
        """封面页：深色背景 + 装饰矩形 + 标题"""
        slide = prs.slides.add_slide(blank)

        # 整体背景
        bg_rect = slide.shapes.add_shape(1, 0, 0, W, H)
        _rgb_fill(bg_rect, C_DARK)

        # 左侧竖向装饰条（亮蓝色）
        accent_bar = slide.shapes.add_shape(1, 0, 0, Inches(0.18), H)
        _rgb_fill(accent_bar, C_ACCENT)

        # 右下角装饰矩形（半透明效果用叠加浅色）
        deco = slide.shapes.add_shape(1, Inches(9.5), Inches(5.2), Inches(3.83), Inches(2.3))
        deco.fill.solid()
        deco.fill.fore_color.rgb = RGBColor(0x0F, 0x34, 0x60)
        deco.line.fill.background()

        # 主标题
        _add_textbox(slide,
                     left=Inches(0.6), top=Inches(2.2),
                     width=Inches(11.5), height=Inches(1.8),
                     text=title, font_size=40, bold=True,
                     color=C_WHITE, align=PP_ALIGN.LEFT)

        # 分隔线（天蓝矩形条）
        sep = slide.shapes.add_shape(1, Inches(0.6), Inches(4.15), Inches(5), Inches(0.06))
        _rgb_fill(sep, C_ACCENT)

        # 副标题 / 要点
        if subtitle_lines:
            sub_text = "  |  ".join(subtitle_lines[:3])
            _add_textbox(slide,
                         left=Inches(0.6), top=Inches(4.4),
                         width=Inches(11.5), height=Inches(0.8),
                         text=sub_text, font_size=16, bold=False,
                         color=C_LIGHT, align=PP_ALIGN.LEFT)

        # 右下角小字
        _add_textbox(slide,
                     left=Inches(10.5), top=Inches(6.9),
                     width=Inches(2.5), height=Inches(0.4),
                     text="ResearchPilot · AI科研助手", font_size=9,
                     color=C_GRAY, align=PP_ALIGN.RIGHT)

    def _add_content_slide(title: str, bullets: list[str], slide_num: int, total: int):
        """内容页：浅灰蓝背景 + 深色标题栏 + 天蓝要点序号"""
        slide = prs.slides.add_slide(blank)

        # 背景
        bg = slide.shapes.add_shape(1, 0, 0, W, H)
        _rgb_fill(bg, C_BG)

        # 顶部标题栏（深蓝矩形，全宽，1.15 英寸高）
        header = slide.shapes.add_shape(1, 0, 0, W, Inches(1.15))
        _rgb_fill(header, C_NAVY)

        # 标题栏左侧天蓝装饰线
        bar = slide.shapes.add_shape(1, 0, 0, Inches(0.12), Inches(1.15))
        _rgb_fill(bar, C_ACCENT)

        # 标题文字
        _add_textbox(slide,
                     left=Inches(0.3), top=Inches(0.15),
                     width=Inches(12.0), height=Inches(0.85),
                     text=title, font_size=24, bold=True,
                     color=C_WHITE, align=PP_ALIGN.LEFT)

        # 内容区背景白色卡片
        card = slide.shapes.add_shape(1, Inches(0.3), Inches(1.3),
                                      Inches(12.73), Inches(5.8))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = RGBColor(0xDC, 0xE8, 0xF5)
        card.line.width = Pt(1)

        # 要点文字
        if bullets:
            tb = slide.shapes.add_textbox(Inches(0.55), Inches(1.55),
                                           Inches(12.2), Inches(5.4))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                if not bullet.strip():
                    continue
                bp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                # 序号用天蓝色小圆点 + 正文
                run_num = bp.add_run()
                run_num.text = f"  {i + 1}.  "
                run_num.font.size = Pt(15)
                run_num.font.bold = True
                run_num.font.color.rgb = C_ACCENT

                run_text = bp.add_run()
                run_text.text = bullet.strip()
                run_text.font.size = Pt(15)
                run_text.font.color.rgb = C_TEXT

                bp.space_after = Pt(8)

        # 底部分隔线
        bot_line = slide.shapes.add_shape(1, Inches(0.3), Inches(7.1),
                                          Inches(12.73), Inches(0.04))
        _rgb_fill(bot_line, C_ACCENT)

        # 页码
        _add_textbox(slide,
                     left=Inches(11.8), top=Inches(7.1),
                     width=Inches(1.3), height=Inches(0.35),
                     text=f"{slide_num} / {total}", font_size=9,
                     color=C_GRAY, align=PP_ALIGN.RIGHT)

        # 左下角品牌小字
        _add_textbox(slide,
                     left=Inches(0.3), top=Inches(7.1),
                     width=Inches(4), height=Inches(0.35),
                     text="ResearchPilot · AI科研助手", font_size=9,
                     color=C_GRAY)

    # ── 解析 AI 输出 ──────────────────────────────────────────────────────────
    slides_raw   = re.split(r'(?i)slide\s*\d+\s*[|｜]', ai_output)
    title_matches = re.findall(r'(?i)slide\s*\d+\s*[|｜]\s*(.+)', ai_output)

    if not title_matches:
        slides_raw   = re.split(r'\n#{1,3}\s+', ai_output)
        title_matches = re.findall(r'\n#{1,3}\s+(.+)', "\n" + ai_output)

    parsed = []
    for title_line, content in zip(title_matches, slides_raw[1:]):
        title = title_line.strip()
        lines = [
            l.lstrip("-•*·▶ ").strip()
            for l in content.strip().splitlines()
            if l.strip() and not re.match(r'(?i)slide\s*\d', l)
        ]
        bullets = [l for l in lines if l][:7]
        parsed.append((title, bullets))

    if not parsed:
        parsed = [("演示文稿", [ai_output[:200]])]

    total = len(parsed)
    for idx, (title, bullets) in enumerate(parsed):
        if idx == 0:
            _add_title_slide(title, bullets)
        else:
            _add_content_slide(title, bullets, idx, total - 1)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_result_paragraph(method: str, results: str, context: str = "") -> str:
    """生成论文结果段落。"""
    user_prompt = RESULT_PARAGRAPH_USER.format(
        method=method, results=results, context=context or "无"
    )
    return chat(EXPRESSION_SYSTEM, user_prompt)


def generate_chart_description(chart_type: str, description: str, findings: str) -> str:
    """生成图表说明文字。"""
    user_prompt = CHART_DESCRIPTION_USER.format(
        chart_type=chart_type, description=description, findings=findings
    )
    return chat(EXPRESSION_SYSTEM, user_prompt)


def generate_model_diagram(description: str) -> bytes:
    """根据模型描述生成科研模型图，返回图片字节流。

    流程：
    1. 调用文本 LLM 将中文描述转化为适合生图的英文 prompt
    2. 调用 DashScope 通义万象原生 SDK 生成图片（自动处理异步轮询）
    3. 下载并返回图片字节
    """
    from http import HTTPStatus
    from dashscope import ImageSynthesis
    from config import LLM_API_KEY

    # Step 1: 生成英文图像 prompt
    prompt_text = chat(
        MODEL_DIAGRAM_SYSTEM,
        MODEL_DIAGRAM_PROMPT_USER.format(description=description),
    )
    prompt_text = prompt_text.strip()

    # Step 2: 调用通义万象原生 SDK（同步阻塞，内部自动处理异步轮询）
    rsp = ImageSynthesis.call(
        api_key=LLM_API_KEY,
        model="wanx2.1-t2i-turbo",
        prompt=prompt_text,
        n=1,
        size="1024*1024",
    )

    if rsp.status_code != HTTPStatus.OK:
        raise RuntimeError(
            f"图像生成失败 (code={rsp.status_code}): {rsp.message}"
        )

    image_url = rsp.output.results[0].url

    # Step 3: 下载图片字节
    with urllib.request.urlopen(image_url) as resp:  # noqa: S310
        image_bytes = resp.read()

    return image_bytes, prompt_text
